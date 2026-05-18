from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "SSH_WEB_FIREWALL_Academic_Presentation.pptx"

MODEL_DIRS = {
    "Firewall": ROOT / "Firewall",
    "Web": ROOT / "WEB LOGS MODEL",
    "SSH": ROOT / "SSH",
}

COLORS = {
    "navy": RGBColor(30, 41, 59),
    "blue": RGBColor(37, 99, 235),
    "green": RGBColor(15, 118, 110),
    "orange": RGBColor(249, 115, 22),
    "red": RGBColor(220, 38, 38),
    "slate": RGBColor(71, 85, 105),
    "light": RGBColor(248, 250, 252),
    "line": RGBColor(203, 213, 225),
    "white": RGBColor(255, 255, 255),
}


def load_metrics() -> dict[str, dict]:
    metrics = {}
    for name, root in MODEL_DIRS.items():
        path = root / "results" / "performance_summary.json"
        with path.open() as handle:
            metrics[name] = json.load(handle)
    return metrics


def add_textbox(slide, x, y, w, h, text, *, size=18, bold=False, color="navy", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    if align:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_bullets(slide, x, y, w, h, bullets, *, size=18, color="slate"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    for index, text in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = text
        paragraph.level = 0
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = COLORS[color]
        paragraph.space_after = Pt(8)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.6, 0.35, 12.2, 0.45, title, size=28, bold=True, color="navy")
    if subtitle:
        add_textbox(slide, 0.62, 0.86, 11.8, 0.35, subtitle, size=13, color="slate")
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.22), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.color.rgb = COLORS["line"]


def add_metric_table(slide, metrics):
    rows = 4
    cols = 6
    table = slide.shapes.add_table(rows, cols, Inches(0.65), Inches(1.55), Inches(12.0), Inches(2.1)).table
    headers = ["Model", "Rows", "Anomalies", "Pred. anomalies", "ROC-AUC", "PR-AUC"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["navy"]
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = COLORS["white"]
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    for row, name in enumerate(["Firewall", "Web", "SSH"], start=1):
        item = metrics[name]
        values = [
            name,
            f"{item['rows']:,}",
            f"{item['anomaly_rows']:,}",
            f"{item['predicted_anomalies']:,}",
            f"{item['roc_auc']:.4f}",
            f"{item['pr_auc']:.4f}",
        ]
        for col, value in enumerate(values):
            cell = table.cell(row, col)
            cell.text = value
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = COLORS["navy"]


def add_picture_or_note(slide, path: Path, x, y, w, h):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        add_textbox(slide, x, y + h / 2 - 0.2, w, 0.4, f"Missing figure:\n{path}", size=12, color="red")


def add_model_result_slide(prs, model_name: str, metrics: dict, figure_names: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    item = metrics[model_name]
    add_title(
        slide,
        f"{model_name} Results",
        f"Rows: {item['rows']:,} | ROC-AUC: {item['roc_auc']:.4f} | PR-AUC: {item['pr_auc']:.4f}",
    )

    fig_dir = MODEL_DIRS[model_name] / "results" / "figures"
    x_positions = [0.65, 4.68, 8.7]
    for x, figure_name in zip(x_positions, figure_names):
        add_picture_or_note(slide, fig_dir / figure_name, x, 1.55, 3.7, 2.95)
        add_textbox(slide, x, 4.55, 3.7, 0.25, figure_name.replace("_", " ").replace(".png", ""), size=10, color="slate", align=PP_ALIGN.CENTER)

    notes = {
        "Firewall": [
            "Ground-truth evaluation uses the Action field.",
            "Strong PR-AUC indicates useful anomaly ranking.",
            "Features capture traffic volume, rate, and port ratios.",
        ],
        "Web": [
            "Evaluation uses error-rate pseudo labels.",
            "PR-AUC is constrained by severe class imbalance.",
            "Hourly host windows make behavior trends visible.",
        ],
        "SSH": [
            "Evaluation uses behavior-rule pseudo labels.",
            "10-minute source-IP windows capture brute-force patterns.",
            "Feature correlations expose suspicious login dynamics.",
        ],
    }
    add_bullets(slide, 0.85, 5.25, 11.7, 1.2, notes[model_name], size=14)


def build() -> None:
    metrics = load_metrics()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.7, 1.0, 12.0, 0.7, "SSH/Web/Firewall Anomaly Detection", size=34, bold=True)
    add_textbox(slide, 0.72, 1.82, 11.6, 0.4, "Multi-source Isolation Forest modeling with live Docker ELK monitoring", size=18, color="slate")
    add_bullets(
        slide,
        0.9,
        3.0,
        11.5,
        2.0,
        [
            "Three security log domains: Firewall, Web, and SSH",
            "Offline training and evaluation with reusable artifacts",
            "Live log generation, file listening, scoring, and Kibana visualization",
        ],
        size=20,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Problem Statement")
    add_bullets(
        slide,
        0.9,
        1.55,
        11.8,
        4.8,
        [
            "Security operations require correlating heterogeneous logs from different systems.",
            "Public datasets often have limited labels or labels that only partially describe attack behavior.",
            "Static model metrics are not enough; operators need live visibility into raw events and predictions.",
            "The project proposes one reproducible architecture for offline validation and live monitoring.",
        ],
        size=20,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Research Objectives")
    add_bullets(
        slide,
        0.9,
        1.55,
        11.8,
        4.8,
        [
            "Design comparable pipelines for Firewall, Web, and SSH anomaly detection.",
            "Engineer domain-specific features while preserving a shared folder architecture.",
            "Train unsupervised Isolation Forest models using normal baselines or available windows.",
            "Evaluate with ROC, precision-recall, confusion matrix, feature correlation, and latency figures.",
            "Integrate Docker-based Elasticsearch, Logstash, and Kibana for live operational dashboards.",
        ],
        size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "End-To-End Architecture")
    stages = [
        "Data collection",
        "Parsing",
        "Feature engineering",
        "Scaling",
        "Isolation Forest",
        "Evaluation",
        "Live scoring",
        "Docker ELK",
    ]
    for idx, stage in enumerate(stages):
        x = 0.55 + (idx % 4) * 3.15
        y = 1.65 + (idx // 4) * 2.15
        shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(2.65), Inches(0.75))
        shape.fill.solid()
        shape.fill.fore_color.rgb = [COLORS["blue"], COLORS["green"], COLORS["orange"], COLORS["slate"]][idx % 4]
        shape.line.color.rgb = COLORS["white"]
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.text = stage
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.runs[0].font.size = Pt(15)
        paragraph.runs[0].font.bold = True
        paragraph.runs[0].font.color.rgb = COLORS["white"]
    add_textbox(slide, 0.9, 6.0, 11.5, 0.5, "The same architecture is reused across all three model folders.", size=17, color="slate", align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Datasets")
    add_bullets(
        slide,
        0.8,
        1.45,
        12.0,
        4.7,
        [
            "Firewall: CSV records with action labels; Action == allow is the normal baseline.",
            "Web: NASA HTTP access log archive; host/hour behavior windows are engineered from Apache-style lines.",
            "SSH: SimpleWeb / University of Twente SSH Dataset 1; auth and Kippo logs are parsed into source-IP windows.",
            "Web and SSH evaluations use pseudo labels because the public logs are not fully attack-labeled.",
        ],
        size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Feature Engineering")
    add_bullets(
        slide,
        0.8,
        1.45,
        12.0,
        4.8,
        [
            "Firewall: bytes per packet, packet rate, byte rate, and port diversity ratio from each CSV row.",
            "Web: request count, unique URLs, bytes, status-code counts, error rate, POST rate, and URL diversity per host/hour.",
            "SSH: failed logins, invalid users, accepted logins, command counts, suspicious commands, failure rates, and burst score per source-IP/10-minute window.",
            "Feature/score correlation figures provide a compact explanation of which engineered features move with anomaly score.",
        ],
        size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Modeling Method")
    add_bullets(
        slide,
        0.8,
        1.45,
        12.0,
        4.8,
        [
            "All models use Isolation Forest, a tree-based unsupervised anomaly detection method.",
            "Numeric features are normalized with StandardScaler before training and scoring.",
            "Firewall trains on Action == allow rows.",
            "SSH trains on low-activity windows.",
            "Web trains on host/hour behavior windows from the NASA HTTP log.",
            "Prediction 1 means normal; prediction -1 means anomalous.",
        ],
        size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Evaluation Protocol")
    add_bullets(
        slide,
        0.8,
        1.45,
        12.0,
        4.8,
        [
            "Firewall: held-out allow rows plus non-allow rows labeled as anomalies.",
            "Web: pseudo anomalies are windows where error_rate > 0.5.",
            "SSH: pseudo anomalies are windows with brute force, enumeration, connection bursts, suspicious commands, or honeypot activity.",
            "Metrics include confusion matrix, ROC-AUC, PR-AUC, score distribution, label balance, feature/score correlation, and inference latency.",
        ],
        size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Results Summary")
    add_metric_table(slide, metrics)
    add_bullets(
        slide,
        0.85,
        4.1,
        11.7,
        1.7,
        [
            "Firewall and SSH show strong ROC-AUC and PR-AUC under their validation definitions.",
            "Web ROC-AUC is high, but PR-AUC is lower because pseudo anomalies are extremely rare.",
            "Latency benchmarks show the models can score at live-demo speed on a local machine.",
        ],
        size=16,
    )

    add_model_result_slide(prs, "Firewall", metrics, ["confusion_matrix.png", "roc_curve.png", "feature_score_correlation.png"])
    add_model_result_slide(prs, "Web", metrics, ["confusion_matrix.png", "precision_recall_curve.png", "feature_score_correlation.png"])
    add_model_result_slide(prs, "SSH", metrics, ["confusion_matrix.png", "roc_curve.png", "feature_score_correlation.png"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Live Scoring")
    add_bullets(
        slide,
        0.8,
        1.45,
        12.0,
        4.7,
        [
            "The generator writes synthetic Firewall, Web, and SSH logs into local live log files.",
            "The listener tails those files, applies the matching feature pipeline, and writes JSONL prediction files.",
            "Firewall is row-based; Web and SSH are window-based but emit a fresh score after each accepted line updates its window.",
            "Each prediction includes model, scored_at, prediction, is_anomaly, anomaly_score, risk_score, and status fields.",
        ],
        size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Kibana Dashboard")
    add_bullets(
        slide,
        0.8,
        1.45,
        12.0,
        4.8,
        [
            "Docker Compose runs Elasticsearch, Kibana, and Logstash.",
            "Logstash tails raw logs and score JSONL files through bind mounts.",
            "Kibana panels show prediction trends, anomaly mix, average risk, raw log volume, firewall actions, web status outcomes, SSH login status, and top SSH source-IP risk.",
            "Normalized status_label fields make Web and SSH visualizations easier to read.",
        ],
        size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Limitations")
    add_bullets(
        slide,
        0.8,
        1.45,
        12.0,
        4.8,
        [
            "Web and SSH evaluations rely on pseudo labels rather than complete expert labels.",
            "Isolation Forest detects unusual behavior; it does not classify attack families.",
            "The live generator is synthetic and intended for integration testing and demonstration.",
            "Feature/score correlation is interpretability support, not causal explanation.",
        ],
        size=18,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Reproducibility")
    add_bullets(
        slide,
        0.8,
        1.45,
        12.0,
        4.8,
        [
            "Create a Python virtual environment and install requirements.txt.",
            "Run each model's preprocess, train, and evaluate scripts from the repository root.",
            "Run evaluation/evaluate_all_models.py to regenerate metrics and figures.",
            "Start ELK with elk/start_elk.sh and install dashboards with elk/setup_kibana.sh.",
            "Run live/listen_and_score.py and live/generate_logs.py in separate terminals.",
        ],
        size=17,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Conclusion")
    add_bullets(
        slide,
        0.8,
        1.45,
        12.0,
        4.8,
        [
            "The project provides a repeatable multi-source anomaly-detection baseline.",
            "The folder architecture makes model-specific pipelines consistent and maintainable.",
            "Evaluation artifacts support offline analysis and academic reporting.",
            "Docker ELK integration turns local live files into useful operational dashboards.",
        ],
        size=19,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    build()
